from pathlib import Path
import re

from django.core.files.base import ContentFile
from django.core.management.base import BaseCommand, CommandError
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from bridges.models import Bridge, Defect, DefectImage

KEY_MAP = {
    "اسم العيب": "title",
    "اسم العنصر": "title",
    "مسمى العنصر": "title",
    "العنصر": "title",
    "رقم العنصر": "item_number",
    "رقم العيب": "code",
    "رقم الحالة": "condition_code",
    "الحالة": "status",
    "شرح الحالة": "description",
    "وصف العيب": "description",
    "العنصر المتضرر": "location",
    "موقع العيب": "location",
    "الطول (م)": "length_m",
    "العرض (م)": "width_m",
    "المساحة (م2)": "area_m2",
    "Count": "count",
    "عدد": "count",
    "Length (m)": "length_m",
    "Width (m)": "width_m",
    "Area (m2)": "area_m2",
    "Area (m²)": "area_m2",
    "Qty. CS1": "count",
    "Qty. CS2": "count",
    "Qty. CS3": "count",
    "Qty. CS4": "count",
    "عيوب": "status",
}

STATUS_KEYWORDS = {
    "open": ["مفتوح", "open"],
    "in_progress": ["جاري", "قيد"],
    "resolved": ["منتهى", "منتهي", "resolved", "closed", "resol"],
}


def normalize_key(value: str) -> str:
    value = value.strip().replace("\xa0", " ")
    lower = value.lower()
    for key, canonical in KEY_MAP.items():
        if key.lower() == lower or key.lower() in lower:
            return canonical
    return lower


def parse_number(value):
    if value is None:
        return None
    text = str(value).strip().replace("،", ".")
    text = re.sub(r"[^0-9.\-]", "", text)
    if not text:
        return None
    try:
        return float(text)
    except ValueError:
        return None


def map_defect_status(value):
    if not value:
        return "open"
    lower = str(value).lower()
    for status, keywords in STATUS_KEYWORDS.items():
        for keyword in keywords:
            if keyword in lower:
                return status
    return "open"


def extract_text_from_shape(shape):
    if hasattr(shape, "has_text_frame") and shape.has_text_frame:
        return shape.text.strip()
    if hasattr(shape, "shapes"):
        texts = []
        for child in shape.shapes:
            child_text = extract_text_from_shape(child)
            if child_text:
                texts.append(child_text)
        return "\n".join(texts).strip()
    return ""


def iter_all_shapes(shape):
    if hasattr(shape, "shapes"):
        for child in shape.shapes:
            yield from iter_all_shapes(child)
    else:
        yield shape


def extract_bridge_name(prs: Presentation, file_name: str) -> str:
    candidate = None
    blacklist = ["موقع على الجسر", "الإحداثيات", "Bridge Number", "Contract Number", "Region"]
    for slide in list(prs.slides)[:3]:
        for shape in iter_all_shapes(slide):
            text = extract_text_from_shape(shape)
            if not text:
                continue
            lines = [line.strip() for line in text.splitlines() if line.strip()]
            if not lines:
                continue
            for line in lines:
                if "جسر" not in line:
                    continue
                if any(bad in line for bad in blacklist):
                    continue
                if len(line) > 5:
                    return line
                if candidate is None:
                    candidate = line
    return candidate or Path(file_name).stem


def extract_coordinates(prs: Presentation):
    coord_re = re.compile(r"X\s*:?[\s]*([0-9.]+)[^0-9a-zA-Z\-]*Y\s*:?[\s]*([0-9.]+)", re.IGNORECASE)
    for slide in list(prs.slides)[:5]:
        for shape in iter_all_shapes(slide):
            text = extract_text_from_shape(shape)
            if not text:
                continue
            match = coord_re.search(text)
            if match:
                return float(match.group(1)), float(match.group(2))
    for slide in list(prs.slides)[:5]:
        for shape in iter_all_shapes(slide):
            if shape.has_table:
                rows = [[cell.text.strip() for cell in row.cells] for row in shape.table.rows]
                data = {normalize_key(r[1]): r[0] for r in rows if len(r) >= 2}
                if "خط العرض" in data and "خط الطول" in data:
                    lat = parse_number(data.get("خط العرض"))
                    lng = parse_number(data.get("خط الطول"))
                    if lat is not None and lng is not None:
                        return lat, lng
    return None


def extract_slide_images(slide):
    images = []
    for shape in iter_all_shapes(slide):
        if hasattr(shape, "image") and shape.image is not None:
            image = shape.image
            images.append({
                "blob": image.blob,
                "ext": image.ext,
                "name": getattr(shape, "name", "image"),
            })
    return images


def parse_defect_tables(prs):
    defects = []
    for slide in prs.slides:
        slide_images = extract_slide_images(slide)
        image_iter = iter(slide_images)
        for shape in iter_all_shapes(slide):
            if not shape.has_table:
                continue
            rows = [[cell.text.strip() for cell in row.cells] for row in shape.table.rows]
            if not any(any(keyword in cell for keyword in ["الحالة", "شرح الحالة", "اسم العيب", "مسمى العنصر", "اسم العنصر", "العنصر المتضرر", "رقم العيب", "عيوب"] for cell in row) for row in rows):
                continue
            row_data = {}
            for row in rows:
                if len(row) >= 2:
                    key = normalize_key(row[0])
                    value = row[1]
                    if key:
                        row_data[key] = value
            if not row_data:
                continue
            status_text = row_data.get("status") or row_data.get("عيوب") or row_data.get("title") or row_data.get("item_number") or "open"
            defects.append({
                "title": status_text,
                "status": map_defect_status(status_text),
                "description": row_data.get("description") or "",
                "location": row_data.get("location") or "",
                "length_m": parse_number(row_data.get("length_m")),
                "width_m": parse_number(row_data.get("width_m")),
                "area_m2": parse_number(row_data.get("area_m2")),
                "count": int(parse_number(row_data.get("count")) or 0),
                "slide_index": slide.slide_id,
                "slide_image": next(image_iter, None),
            })
    return defects


class Command(BaseCommand):
    help = "استيراد الجسور والعيوب من ملفات PPTX في مجلد data إلى قاعدة البيانات"

    def add_arguments(self, parser):
        parser.add_argument(
            "--data-dir",
            default="data",
            help="مسار المجلد الذي يحتوي على ملفات PPTX",
        )
        parser.add_argument(
            "--file",
            help="اسم ملف PPTX واحد للاستيراد من مجلد البيانات",
        )
        parser.add_argument(
            "--clear",
            action="store_true",
            help="حذف البيانات الحالية قبل الاستيراد",
        )

    def handle(self, *args, **options):
        data_dir = Path(options["data_dir"])
        if not data_dir.exists() or not data_dir.is_dir():
            raise CommandError(f"مجلد البيانات غير موجود: {data_dir}")

        if options["clear"]:
            self.stdout.write(self.style.WARNING("حذف جميع الجسور والعيوب والصور الحالية..."))
            DefectImage.objects.all().delete()
            Defect.objects.all().delete()
            Bridge.objects.all().delete()

        if options.get("file"):
            file_path = data_dir / options["file"]
            if not file_path.exists() or not file_path.is_file():
                raise CommandError(f"لم يتم العثور على الملف: {file_path}")
            files = [file_path]
        else:
            files = [f for f in sorted(data_dir.glob("*.pptx")) if not f.name.startswith("~$")]
            if not files:
                raise CommandError(f"لم يتم العثور على ملفات PPTX في {data_dir}")

        imported_bridges = 0
        imported_defects = 0
        imported_images = 0

        for path in files:
            self.stdout.write(f"استيراد الملف: {path.name}")
            prs = Presentation(path)
            bridge_name = extract_bridge_name(prs, path.name)
            coords = extract_coordinates(prs)
            defaults = {
                "status": "maintenance",
                "lat": coords[0] if coords else 0.0,
                "lng": coords[1] if coords else 0.0,
            }
            bridge, created = Bridge.objects.get_or_create(name=bridge_name, defaults=defaults)
            if created:
                imported_bridges += 1
                self.stdout.write(self.style.SUCCESS(f"  أنشئ جسر جديد: {bridge_name}"))
            else:
                if coords:
                    bridge.lat, bridge.lng = coords
                    bridge.save(update_fields=["lat", "lng"])
                self.stdout.write(f"  تم العثور على جسر موجود: {bridge_name}")

            defects = parse_defect_tables(prs)
            for defect_data in defects:
                title = defect_data["title"]
                description = defect_data["description"]
                status = defect_data["status"]
                length_m = defect_data["length_m"] or 0.0
                width_m = defect_data["width_m"] or 0.0
                area_m2 = defect_data["area_m2"] or 0.0
                count = defect_data["count"]
                defect = Defect.objects.create(
                    bridge=bridge,
                    title=title,
                    status=status,
                    description=description,
                    length_m=length_m,
                    width_m=width_m,
                    area_m2=area_m2,
                    count=count,
                )
                imported_defects += 1
                image_info = defect_data.get("slide_image")
                if image_info:
                    file_name = f"pptx_{bridge.id}_{defect.id}.{image_info['ext']}"
                    di = DefectImage(defect=defect, caption=title)
                    di.image.save(file_name, ContentFile(image_info["blob"]), save=True)
                    imported_images += 1
            self.stdout.write(f"  أضيفت {len(defects)} عيوب إلى الجسر")

        self.stdout.write(self.style.SUCCESS(
            f"انتهى الاستيراد: {imported_bridges} جسور، {imported_defects} عيوب، {imported_images} صور."))
